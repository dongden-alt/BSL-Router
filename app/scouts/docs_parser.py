"""
Scout.docs_parser — Document Intelligence Scout

Intercepts incoming ChatCompletionRequests containing attached documents.
Extracts text from PDF, DOCX, XLSX, PPTX, and TXT files.
If the extracted text exceeds the configured threshold (default 8,000 tokens),
it asks a cheap model to summarize the document contextually based on the user's intent,
before passing the context to the primary model.
"""

import httpx
import base64
import io
import asyncio
from app.models import ChatCompletionRequest, MessageContentPart
from app.utils.url_normalization import build_custom_text_upstream_url

# Optional imports for document parsing
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def _approximate_tokens(text: str) -> int:
    return len(text) // 4


def _extract_intent(request: ChatCompletionRequest) -> str:
    """Find the last user message to use as the intent context for summarization."""
    for msg in reversed(request.messages):
        if msg.role == "user":
            content = msg.content
            if isinstance(content, str):
                return content
            elif isinstance(content, list):
                text_parts = []
                for p in content:
                    if isinstance(p, dict) and p.get("type") == "text":
                        text_parts.append(p.get("text", ""))
                    elif isinstance(p, MessageContentPart) and p.type == "text" and p.text:
                        text_parts.append(p.text)
                return " ".join(text_parts)
    return "Analyze and extract all important information from this document."


def _parse_pdf(data: bytes) -> str:
    if not fitz:
        return "[PDF Parsing Failed: PyMuPDF not installed]"
    doc = fitz.open(stream=data, filetype="pdf")
    text = []
    try:
        for page in doc:
            text.append(page.get_text())
    finally:
        doc.close()
    return "\n".join(text)


def _parse_docx(data: bytes) -> str:
    if not docx:
        return "[DOCX Parsing Failed: python-docx not installed]"
    doc = docx.Document(io.BytesIO(data))
    return "\n".join([p.text for p in doc.paragraphs])


def _parse_xlsx(data: bytes) -> str:
    if not openpyxl:
        return "[XLSX Parsing Failed: openpyxl not installed]"
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    text = []
    for sheet in wb.worksheets:
        text.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(c) if c is not None else "" for c in row])
            if row_text.strip():
                text.append(row_text)
    return "\n".join(text)


def _parse_pptx(data: bytes) -> str:
    if not Presentation:
        return "[PPTX Parsing Failed: python-pptx not installed]"
    prs = Presentation(io.BytesIO(data))
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def _parse_document_sync(mime_type: str, b64_data: str) -> str:
    """CPU-bound parsing logic to be run in a thread pool."""
    try:
        data = base64.b64decode(b64_data)
    except Exception as e:
        return f"[Failed to decode base64 document: {e}]"

    try:
        if "pdf" in mime_type:
            return _parse_pdf(data)
        elif "wordprocessingml.document" in mime_type or "msword" in mime_type or mime_type == "docx":
            return _parse_docx(data)
        elif "spreadsheetml.sheet" in mime_type or "ms-excel" in mime_type or mime_type == "xlsx":
            return _parse_xlsx(data)
        elif "presentationml.presentation" in mime_type or "ms-powerpoint" in mime_type or mime_type == "pptx":
            return _parse_pptx(data)
        elif "text/plain" in mime_type or mime_type in ["txt", "md", "csv"]:
            return data.decode("utf-8", errors="ignore")
        else:
            return f"[Unsupported document type: {mime_type}]"
    except Exception as e:
        return f"[Parsing Error ({mime_type}): {e}]"


async def _summarize_document(
    text: str,
    intent: str,
    http_client: httpx.AsyncClient,
    conn: dict,
    model: str
) -> str:
    """Call the cheap model to extract only relevant parts."""
    prompt = (
        f"You are a document intelligence assistant. The user's current intent/request is:\n"
        f"\"\"\"{intent}\"\"\"\n\n"
        f"Based ONLY on the user's intent, extract the relevant information from the following document. "
        f"If the document is extremely large, summarize the key points that answer the user's query. "
        f"Do not hallucinate outside information.\n\n"
        f"--- DOCUMENT START ---\n{text}\n--- DOCUMENT END ---"
    )

    # Wire format travels with the connection (model_resolver injects the
    # provider-level `format`). This scout previously had no anthropic branch
    # at all, so an anthropic-format summarization model received an OpenAI
    # body with Bearer auth and would reject the request.
    fmt = str(conn.get("format") or "openai").lower()
    is_anthropic_fmt = fmt == "anthropic"

    # URL construction is delegated to the same builder main.py uses, so the
    # scouts cannot drift from the primary routing path.
    endpoint = build_custom_text_upstream_url(
        conn.get("base_url", ""),
        "anthropic" if is_anthropic_fmt else "openai",
    )

    if is_anthropic_fmt:
        payload = {
            "model": model,
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 8192,
        }
        headers = {
            "x-api-key": conn.get("api_key", ""),
            "anthropic-version": "2023-06-01",
            "Content-Type": "application/json",
        }
        resp = await http_client.post(endpoint, json=payload, headers=headers, timeout=600.0)
        resp.raise_for_status()
        blocks = resp.json().get("content", [])
        return " ".join(
            b.get("text", "") for b in blocks
            if isinstance(b, dict) and b.get("type") == "text"
        )

    payload = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": 8192,
        "stream": False
    }

    headers = {
        "Authorization": f"Bearer {conn['api_key']}",
        "Content-Type": "application/json"
    }

    resp = await http_client.post(endpoint, json=payload, headers=headers, timeout=600.0)
    resp.raise_for_status()

    data = resp.json()
    return data.get("choices", [{}])[0].get("message", {}).get("content", "")


def _extract_doc_info(part) -> tuple[str, str]:
    """Returns (mime_type, base64_data) or (None, None)."""
    # 1. Anthropic native document block
    if isinstance(part, dict) and part.get("type") == "document":
        src = part.get("source", {})
        return src.get("media_type", "unknown"), src.get("data", "")
    elif isinstance(part, MessageContentPart) and part.type == "document" and part.source:
        return part.source.get("media_type", "unknown"), part.source.get("data", "")
        
    # 2. OpenAI image_url hijacked for documents (e.g. data:application/pdf;base64,...)
    url = ""
    if isinstance(part, dict) and part.get("type") == "image_url":
        url = part.get("image_url", {}).get("url", "")
    elif isinstance(part, MessageContentPart) and part.type == "image_url" and part.image_url:
        url = part.image_url.get("url", "")

    if url.startswith("data:"):
        # Format: data:application/pdf;base64,JVBERi0xLj...
        try:
            header, b64 = url.split(",", 1)
            mime = header.replace("data:", "").replace(";base64", "")
            # Only intercept known document mimes, ignore actual images
            if "pdf" in mime or "officedocument" in mime or "text/" in mime:
                return mime, b64
        except Exception:
            pass

    return None, None


async def parse_documents(
    request: ChatCompletionRequest,
    http_client: httpx.AsyncClient,
    config: dict
) -> ChatCompletionRequest:
    """
    Main entry point for the Docs Parser Scout.
    """
    t = config.get("tools", {})
    if not t.get("docs_parser_enabled", True):
        return request

    # Scan for documents
    doc_hits = []
    for msg_idx, msg in enumerate(request.messages):
        if not isinstance(msg.content, list):
            continue
        for part_idx, part in enumerate(msg.content):
            mime, b64 = _extract_doc_info(part)
            if mime and b64:
                doc_hits.append((msg_idx, part_idx, mime, b64))

    if not doc_hits:
        return request

    intent = _extract_intent(request)
    skip_threshold = t.get("docs_skip_threshold", 8000)
    
    # Resolve summarization model connection (supports Combo aliases)
    summary_model = t.get("docs_summary_model", "gpt-4o-mini")
    from app.utils.model_resolver import resolve_model_conn
    active_conn, summary_model = resolve_model_conn(config, summary_model)

    if not active_conn:
        print(f"[DocsParser] No active connection found for summarization model '{t.get('docs_summary_model', 'gpt-4o-mini')}' — documents will not be summarized, only extracted.")

    # Process each document
    async def process_doc(msg_idx, part_idx, mime, b64):
        # 1. Parse using a background thread (CPU-bound)
        raw_text = await asyncio.to_thread(_parse_document_sync, mime, b64)
        
        # 2. Check token threshold
        tokens = _approximate_tokens(raw_text)
        
        # 3. Summarize if needed and we have a valid connection
        if tokens > skip_threshold and active_conn:
            try:
                final_text = await _summarize_document(raw_text, intent, http_client, active_conn, summary_model)
                final_text = f"[Document Summarized by {summary_model}]\n\n{final_text}"
            except Exception as e:
                # Fail-safe: truncate to avoid blowing context window
                excerpt = raw_text[:4000] + "... [truncated — summarization failed]" if len(raw_text) > 4000 else raw_text
                final_text = f"[Summarization Failed — {e}]\n\n{excerpt}"
        else:
            final_text = raw_text
            
        return (msg_idx, part_idx), final_text

    tasks = [process_doc(m, p, mime, b64) for m, p, mime, b64 in doc_hits]
    results = await asyncio.gather(*tasks)
    processed_docs = dict(results)

    # Rewrite message content
    for msg_idx, msg in enumerate(request.messages):
        if not isinstance(msg.content, list):
            continue
        new_parts = []
        for part_idx, part in enumerate(msg.content):
            key = (msg_idx, part_idx)
            if key in processed_docs:
                new_parts.append({
                    "type": "text",
                    "text": f"--- Attached Document ---\n{processed_docs[key]}\n------------------------"
                })
            else:
                new_parts.append(part)
        msg.content = new_parts

    return request
